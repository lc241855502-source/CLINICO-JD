import pdfplumber
from docx import Document
from pptx import Presentation
import io

def parse_document(file_bytes, file_name):
    """解析多种格式文档，返回纯文本内容"""
    ext = file_name.split('.')[-1].lower()
    
    if ext == 'docx':
        return _parse_docx(file_bytes)
    elif ext == 'pdf':
        return _parse_pdf(file_bytes)
    elif ext == 'pptx':
        return _parse_pptx(file_bytes)
    elif ext == 'txt':
        return file_bytes.decode('utf-8', errors='ignore')
    else:
        raise ValueError(f"不支持的文件格式: .{ext}")

def _parse_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    
    # 提取表格内容
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(' | '.join(cells))
    
    return '\n'.join(paragraphs)

def _parse_pdf(file_bytes):
    text = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
    return '\n'.join(text)

def _parse_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
    return '\n'.join(text)
